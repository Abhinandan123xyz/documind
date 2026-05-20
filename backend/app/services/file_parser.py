import fitz  # pymupdf - for PDF
import os
from typing import List, Dict
import traceback

def chunk_text(text: str, source: str, page: int = 1, chunk_size: int = 500, overlap: int = 50) -> List[Dict]:
    """Universal chunker — works for any text content."""
    if not text or not text.strip():
        print(f"⚠️ Empty text for {source} page {page}")
        return []
        
    words = text.split()
    chunks = []
    start = 0
    chunk_index = 0

    while start < len(words):
        end = start + chunk_size
        chunk_words = words[start:end]
        chunk_text_content = " ".join(chunk_words)
        
        if len(chunk_text_content.strip()) > 50:  # Only keep substantial chunks
            chunks.append({
                "text": chunk_text_content,
                "page": page,
                "chunk_index": chunk_index,
                "source": source
            })
            chunk_index += 1
        start += chunk_size - overlap

    print(f"✅ Created {len(chunks)} chunks from {source} page {page}")
    return chunks


def parse_pdf(file_path: str) -> List[Dict]:
    """Parse PDF file into chunks."""
    try:
        doc = fitz.open(file_path)
        chunks = []
        source = os.path.basename(file_path)

        print(f"📄 Parsing PDF: {source} ({len(doc)} pages)")
        
        for page_num in range(len(doc)):
            text = doc[page_num].get_text("text").strip()
            if text:
                chunks.extend(chunk_text(text, source, page_num + 1))

        doc.close()
        print(f"✅ Total chunks from PDF: {len(chunks)}")
        return chunks
        
    except Exception as e:
        print(f"❌ Error parsing PDF {file_path}: {e}")
        traceback.print_exc()
        raise


def parse_docx(file_path: str) -> List[Dict]:
    """Parse DOCX file into chunks."""
    try:
        from docx import Document
        doc = Document(file_path)
        source = os.path.basename(file_path)
        
        print(f"📝 Parsing DOCX: {source}")
        
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        chunks = chunk_text(full_text, source)
        
        print(f"✅ Total chunks from DOCX: {len(chunks)}")
        return chunks
        
    except Exception as e:
        print(f"❌ Error parsing DOCX {file_path}: {e}")
        traceback.print_exc()
        raise


def parse_pptx(file_path: str) -> List[Dict]:
    """Parse PPTX file into chunks."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        source = os.path.basename(file_path)
        chunks = []

        print(f"📊 Parsing PPTX: {source} ({len(prs.slides)} slides)")

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text = "\n".join(slide_text)
                chunks.extend(chunk_text(text, source, slide_num))

        print(f"✅ Total chunks from PPTX: {len(chunks)}")
        return chunks
        
    except Exception as e:
        print(f"❌ Error parsing PPTX {file_path}: {e}")
        traceback.print_exc()
        raise


def parse_excel_csv(file_path: str) -> List[Dict]:
    """Parse Excel/CSV file into chunks."""
    try:
        import pandas as pd
        source = os.path.basename(file_path)

        print(f"📈 Parsing Excel/CSV: {source}")
        
        if file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        chunks = []
        chunk_size = 50  # rows per chunk
        
        for i in range(0, len(df), chunk_size):
            batch = df.iloc[i:i+chunk_size]
            text = batch.to_string(index=False)
            chunks.append({
                "text": text,
                "page": (i // chunk_size) + 1,
                "chunk_index": i // chunk_size,
                "source": source
            })

        print(f"✅ Total chunks from Excel/CSV: {len(chunks)}")
        return chunks
        
    except Exception as e:
        print(f"❌ Error parsing Excel/CSV {file_path}: {e}")
        traceback.print_exc()
        raise


def parse_url(url: str) -> List[Dict]:
    """Scrape and chunk text from a web URL."""
    try:
        import requests
        from bs4 import BeautifulSoup
        import re

        print(f"🌐 Scraping URL: {url}")
        
        headers = {"User-Agent": "Mozilla/5.0 (compatible; DocuMindBot/1.0)"}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")

        # Remove noise
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()

        text = soup.get_text(separator=" ", strip=True)
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text).strip()

        source = url[:80]  # truncate long URLs
        chunks = chunk_text(text, source)
        
        print(f"✅ Total chunks from URL: {len(chunks)}")
        return chunks
        
    except Exception as e:
        print(f"❌ Error scraping URL {url}: {e}")
        traceback.print_exc()
        raise


def parse_file(file_path: str) -> List[Dict]:
    """Auto-detect file type and parse accordingly."""
    ext = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".xlsx": parse_excel_csv,
        ".xls":  parse_excel_csv,
        ".csv":  parse_excel_csv,
    }

    parser = parsers.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")

    return parser(file_path)